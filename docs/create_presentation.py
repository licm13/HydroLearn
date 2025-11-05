"""
PowerPoint Presentation Generator for HydroLearn
HydroLearn 教学演示文稿生成器

This script generates a comprehensive PowerPoint presentation for teaching
hydrological modeling with the HydroLearn models.

这个脚本为HydroLearn水文模型教学生成综合的PowerPoint演示文稿。

Author: HydroLearn Teaching Team
Date: 2024
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    return slide


def create_content_slide(prs, title, content_list):
    """Create a slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide


def create_image_slide(prs, title, image_path, caption=""):
    """Create a slide with an image"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add image if it exists
    if os.path.exists(image_path):
        left = Inches(0.5)
        top = Inches(1.3)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(9))
    
    # Add caption if provided
    if caption:
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = caption
        caption_para.font.size = Pt(14)
        caption_para.alignment = PP_ALIGN.CENTER
    
    return slide


def create_two_column_slide(prs, title, left_content, right_content):
    """Create a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Left column
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.25)
    height = Inches(5)
    
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # Right column
    left = Inches(5.25)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    return slide


def main():
    """
    Generate comprehensive teaching PowerPoint presentation.
    生成综合教学PowerPoint演示文稿。
    """
    print("=" * 80)
    print("Generating HydroLearn Teaching Presentation")
    print("生成 HydroLearn 教学演示文稿")
    print("=" * 80)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Hydrological Modeling with HydroLearn\n水文模型教学库",
        "An Educational Repository for Rainfall-Runoff Modeling\n降雨径流模型教学库"
    )
    
    # Slide 2: Course Objectives
    create_content_slide(
        prs,
        "Learning Objectives / 学习目标",
        [
            "• Understand the fundamental concepts of hydrological modeling",
            "  理解水文模拟的基本概念",
            "• Learn different runoff generation mechanisms",
            "  学习不同的产流机制",
            "• Master parameter calibration and validation techniques",
            "  掌握参数率定和验证技术",
            "• Compare performance of different hydrological models",
            "  比较不同水文模型的性能",
            "• Apply models to real-world problems",
            "  将模型应用于实际问题"
        ]
    )
    
    # Slide 3: Hydrological Cycle
    create_image_slide(
        prs,
        "The Hydrological Water Cycle / 水文循环",
        "../figures/water_cycle_diagram.png",
        "Water Balance: P = ET + Q + ΔS"
    )
    
    # Slide 4: What is a Hydrological Model?
    create_content_slide(
        prs,
        "What is a Hydrological Model? / 什么是水文模型?",
        [
            "• Mathematical representation of the hydrological cycle",
            "  水文循环的数学表示",
            "• Transforms inputs (P, T, PET) into outputs (Q)",
            "  将输入(降水、温度、蒸散发)转换为输出(径流)",
            "• Three main types:",
            "  - Physical-based (physically-based processes)",
            "  - Conceptual (simplified representations)",
            "  - Data-driven (statistical/machine learning)",
            "• We focus on conceptual models in this course",
            "  本课程重点关注概念性模型"
        ]
    )
    
    # Slide 5: Model Comparison Table
    create_image_slide(
        prs,
        "HydroLearn Model Portfolio / 模型组合",
        "../figures/model_comparison_table.png",
        "Six different hydrological models for various applications / 六种不同的水文模型"
    )
    
    # Slide 6: Xinanjiang Model Introduction
    create_two_column_slide(
        prs,
        "Xinanjiang Model (新安江模型)",
        [
            "📍 Origin:",
            "• Developed by Zhao Ren-jun (1973)",
            "• Hohai University, China",
            "",
            "🎯 Key Features:",
            "• Saturation excess mechanism",
            "• Three-layer evapotranspiration",
            "• Parabolic distribution curves",
            "• Best for humid regions",
            "",
            "📊 Parameters: 13",
            "• Evaporation: K, C",
            "• Soil: WM, WUM, WLM, B",
            "• Runoff: SM, EX, IMP",
            "• Routing: KI, KG, CI, CG"
        ],
        [
            "🔧 Main Components:",
            "",
            "1. Three-Layer ET Module",
            "   • Upper layer (WUM)",
            "   • Lower layer (WLM)",
            "   • Deep layer (WDM)",
            "",
            "2. Runoff Generation",
            "   • Saturation excess",
            "   • Parabolic curve",
            "",
            "3. Runoff Separation",
            "   • Surface runoff",
            "   • Interflow",
            "   • Groundwater",
            "",
            "4. Flow Routing",
            "   • Linear reservoirs"
        ]
    )
    
    # Slide 7: Xinanjiang Model Flowchart
    create_image_slide(
        prs,
        "Xinanjiang Model Structure / 新安江模型结构",
        "../figures/xinanjiang_model_flowchart.png",
        "Complete model flowchart showing all components / 完整模型流程图"
    )
    
    # Slide 8: HBV Model Introduction
    create_two_column_slide(
        prs,
        "HBV Model",
        [
            "📍 Origin:",
            "• Swedish Meteorological and",
            "  Hydrological Institute (1970s)",
            "• Widely used worldwide",
            "",
            "🎯 Key Features:",
            "• Temperature-driven processes",
            "• Snow accumulation and melt",
            "• Soil moisture accounting",
            "• Three-component runoff",
            "• Best for cold/temperate regions",
            "",
            "📊 Parameters: 13",
            "• Snow: TT, CFMAX, CWH, CFR",
            "• Soil: FC, LP, BETA",
            "• Response: PERC, UZL, K0, K1, K2",
            "• Routing: MAXBAS"
        ],
        [
            "🔧 Main Modules:",
            "",
            "1. Snow Routine",
            "   • Degree-day method",
            "   • Snow accumulation",
            "   • Snowmelt calculation",
            "",
            "2. Soil Moisture Routine",
            "   • Actual ET calculation",
            "   • Recharge to response zone",
            "",
            "3. Response Routine",
            "   • Upper zone (SUZ)",
            "   • Lower zone (SLZ)",
            "   • Q0 (quick runoff)",
            "   • Q1 (interflow)",
            "   • Q2 (baseflow)",
            "",
            "4. Routing",
            "   • Triangular weighting"
        ]
    )
    
    # Slide 9: HBV Model Flowchart
    create_image_slide(
        prs,
        "HBV Model Structure / HBV模型结构",
        "../figures/hbv_model_flowchart.png",
        "Complete model flowchart with snow, soil, and response routines / 完整流程图"
    )
    
    # Slide 10: Other Models Overview
    create_content_slide(
        prs,
        "Other Models in HydroLearn / 其他模型",
        [
            "🔹 Tank Model (タンクモデル)",
            "  • Multi-reservoir conceptual model from Japan",
            "  • Flexible structure for various catchments",
            "",
            "🔹 GR4J Model",
            "  • Only 4 parameters - very parsimonious!",
            "  • Excellent for daily streamflow simulation",
            "",
            "🔹 Sacramento Model (SAC-SMA)",
            "  • Continuous soil moisture accounting",
            "  • Used in operational flood forecasting",
            "",
            "🔹 SCS-CN + Unit Hydrograph",
            "  • Event-based model for storm analysis",
            "  • Design flood estimation"
        ]
    )
    
    # Slide 11: Model Calibration
    create_content_slide(
        prs,
        "Model Calibration & Validation / 参数率定与验证",
        [
            "📌 What is Calibration?",
            "  • Finding optimal parameter values",
            "  • Maximizing model performance metrics",
            "  • Using historical observed data",
            "",
            "📌 Common Performance Metrics:",
            "  • Nash-Sutcliffe Efficiency (NSE): -∞ to 1 (1 = perfect)",
            "  • Root Mean Square Error (RMSE): lower is better",
            "  • Coefficient of Determination (R²)",
            "",
            "📌 Best Practices:",
            "  • Split data: Calibration (60-70%) + Validation (30-40%)",
            "  • Use independent validation period",
            "  • Check parameter physical meaning"
        ]
    )
    
    # Slide 12: How to Use HydroLearn
    create_content_slide(
        prs,
        "Getting Started with HydroLearn / 开始使用",
        [
            "📥 Step 1: Installation",
            "  git clone https://github.com/licm13/HydroLearn.git",
            "  pip install -r requirements.txt",
            "",
            "🔬 Step 2: Run Examples",
            "  python examples.py          # Compare all models",
            "  python hbv_model.py         # Run HBV model",
            "  python xinanjiang_model.py  # Run Xinanjiang",
            "",
            "📊 Step 3: Explore Visualizations",
            "  Check the figures/ directory for plots",
            "",
            "🎓 Step 4: Try Calibration",
            "  python calibration_example.py  # Parameter optimization"
        ]
    )
    
    # Slide 13: Repository Structure
    create_content_slide(
        prs,
        "Repository Structure / 项目结构",
        [
            "📁 HydroLearn/",
            "  ├── 📄 README.md              # Documentation",
            "  ├── 📄 requirements.txt       # Dependencies",
            "  ├── 📂 data/                  # Sample datasets",
            "  ├── 📂 figures/               # Output visualizations",
            "  ├── 📂 docs/                  # Teaching materials",
            "  ├── 📜 xinanjiang_model.py    # Xinanjiang implementation",
            "  ├── 📜 hbv_model.py           # HBV implementation",
            "  ├── 📜 tank_model.py          # Tank model",
            "  ├── 📜 gr4j_model.py          # GR4J model",
            "  ├── 📜 sacramento_model.py    # Sacramento model",
            "  ├── 📜 event_model_scs_uh.py  # SCS-CN + UH",
            "  ├── 📜 examples.py            # Model comparison",
            "  └── 📜 calibration_example.py # Calibration demo"
        ]
    )
    
    # Slide 14: Key Concepts to Remember
    create_content_slide(
        prs,
        "Key Concepts / 关键概念",
        [
            "💧 Water Balance: P = ET + Q + ΔS",
            "  • Input = Output + Storage Change",
            "",
            "🌊 Runoff Generation Mechanisms:",
            "  • Saturation Excess (Xinanjiang): Runoff when soil is saturated",
            "  • Infiltration Excess (SCS-CN): Runoff when rainfall > infiltration",
            "",
            "📈 Model Complexity Trade-off:",
            "  • More parameters ≠ Better performance",
            "  • Balance between complexity and data availability",
            "",
            "🎯 Model Selection:",
            "  • Choose based on region, data, and purpose",
            "  • No single 'best' model for all situations"
        ]
    )
    
    # Slide 15: Assignments and Practice
    create_content_slide(
        prs,
        "Practice Assignments / 练习作业",
        [
            "📝 Assignment 1: Model Exploration",
            "  • Run all models with default parameters",
            "  • Compare outputs and water balance",
            "  • Analyze strengths and weaknesses",
            "",
            "📝 Assignment 2: Parameter Sensitivity",
            "  • Vary key parameters (±20%)",
            "  • Observe impact on simulated discharge",
            "  • Identify most sensitive parameters",
            "",
            "📝 Assignment 3: Model Calibration",
            "  • Calibrate Xinanjiang or HBV model",
            "  • Validate on independent period",
            "  • Compare NSE, RMSE, and visual fit",
            "",
            "📝 Assignment 4: Report and Presentation",
            "  • Write technical report on findings",
            "  • Present model comparison results"
        ]
    )
    
    # Slide 16: Additional Resources
    create_content_slide(
        prs,
        "Additional Resources / 额外资源",
        [
            "📚 Documentation:",
            "  • README.md - Complete usage guide",
            "  • Model docstrings - Detailed parameter descriptions",
            "  • Jupyter notebooks - Interactive tutorials",
            "",
            "🔗 References:",
            "  • Zhao (1992) - Xinanjiang model original paper",
            "  • Bergström (1992) - HBV model documentation",
            "  • Perrin et al. (2003) - GR4J model",
            "",
            "💻 GitHub Repository:",
            "  • https://github.com/licm13/HydroLearn",
            "  • Issues and discussions welcome",
            "",
            "📧 Contact:",
            "  • Questions? Open an issue on GitHub"
        ]
    )
    
    # Slide 17: Summary
    create_content_slide(
        prs,
        "Summary / 总结",
        [
            "✅ We covered:",
            "  • Fundamentals of hydrological modeling",
            "  • Six different rainfall-runoff models",
            "  • Model structures and components",
            "  • Calibration and validation techniques",
            "  • Practical implementation with Python",
            "",
            "🎯 Key Takeaways:",
            "  • Models are tools - understand their assumptions",
            "  • Calibration requires both data and physical insight",
            "  • Validation is essential for credible results",
            "  • Different models for different purposes",
            "",
            "🚀 Next Steps:",
            "  • Practice with the repository",
            "  • Complete assignments",
            "  • Apply to your own catchment data"
        ]
    )
    
    # Slide 18: Thank You
    create_title_slide(
        prs,
        "Thank You! / 谢谢!",
        "Questions and Discussions Welcome\n欢迎提问和讨论\n\nHydroLearn Teaching Team"
    )
    
    # Save presentation
    output_path = "HydroLearn_Teaching_Presentation.pptx"
    prs.save(output_path)
    
    print(f"\n✓ PowerPoint presentation created successfully!")
    print(f"✓ 成功创建PowerPoint演示文稿!")
    print(f"✓ Saved as: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print("\n" + "=" * 80)
    
    return output_path


if __name__ == "__main__":
    main()

